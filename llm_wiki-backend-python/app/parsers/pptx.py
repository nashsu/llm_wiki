"""PPTX presentation parser using python-pptx."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.parsers.base import DocumentParser, ExtractedImage, ParseResult
from app.parsers.registry import register_parser

logger = logging.getLogger(__name__)


def _extract_shape_text(shape) -> str:
    """Recursively extract text from a slide shape."""
    parts: list[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

    if shape.has_table:
        table = shape.table
        rows = []
        col_count = len(table.columns)
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)
        if rows:
            # Normalise
            for row in rows:
                while len(row) < col_count:
                    row.append("")
            table_lines: list[str] = []
            table_lines.append("| " + " | ".join(rows[0]) + " |")
            table_lines.append("| " + " | ".join(["---"] * col_count) + " |")
            for row in rows[1:]:
                table_lines.append("| " + " | ".join(row) + " |")
            parts.append("\n".join(table_lines))

    # Handle group shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            child_text = _extract_shape_text(child)
            if child_text:
                parts.append(child_text)

    return "\n".join(parts)


@register_parser()
class PPTXParser(DocumentParser):
    """Parser for PowerPoint presentations (``.pptx``).

    Each slide is extracted as a ``## Slide N`` section containing its text
    content.  Embedded images are extracted where possible.
    """

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    @property
    def supports_images(self) -> bool:
        return True

    def parse(self, file_path: str) -> ParseResult:
        path = Path(file_path)
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            logger.error("Failed to open PPTX %s: %s", path, exc)
            return ParseResult(
                text="",
                images=[],
                metadata={
                    "source": str(path),
                    "extension": path.suffix,
                    "parser": self.name,
                },
                success=False,
                error=f"Failed to open PPTX: {exc}",
            )

        try:
            text_parts: list[str] = []
            images: list[ExtractedImage] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    # Extract text
                    shape_text = _extract_shape_text(shape)
                    if shape_text:
                        slide_texts.append(shape_text)

                    # Extract images
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            content_type = image.content_type
                            # Derive a filename
                            ext = content_type.split("/")[-1] if "/" in content_type else "bin"
                            if ext == "jpeg":
                                ext = "jpg"
                            filename = f"slide{slide_num}_{shape.name or 'image'}.{ext}"

                            images.append(
                                ExtractedImage(
                                    filename=filename,
                                    data=image_bytes,
                                    mime_type=content_type,
                                )
                            )
                        except Exception as img_exc:
                            logger.warning(
                                "Failed to extract image from slide %d: %s",
                                slide_num,
                                img_exc,
                            )

                slide_content = "\n\n".join(slide_texts)
                if slide_content.strip():
                    text_parts.append(f"## Slide {slide_num}\n\n{slide_content}")

            text = "\n\n".join(text_parts)

            metadata: dict = {
                "source": str(path),
                "extension": path.suffix,
                "parser": self.name,
                "slide_count": len(prs.slides),
            }

            return ParseResult(text=text, images=images, metadata=metadata)

        except Exception as exc:
            logger.error("Error parsing PPTX %s: %s", path, exc)
            return ParseResult(
                text="",
                images=[],
                metadata={
                    "source": str(path),
                    "extension": path.suffix,
                    "parser": self.name,
                },
                success=False,
                error=f"Error parsing PPTX: {exc}",
            )
